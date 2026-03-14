import json
from pathlib import Path
from datetime import datetime
from typing import Optional

try:
    from .models import AssessmentData
except ImportError:
    from models import AssessmentData


MATURITY_BANDS = [
    {"min": 1.0, "max": 1.8, "label": "Initial", "color": "#ef4444"},
    {"min": 1.8, "max": 2.6, "label": "Basic", "color": "#f97316"},
    {"min": 2.6, "max": 3.4, "label": "Intermediate", "color": "#eab308"},
    {"min": 3.4, "max": 4.2, "label": "Advanced", "color": "#84cc16"},
    {"min": 4.2, "max": 5.0, "label": "Optimizing", "color": "#22c55e"},
]

SCORE_LABELS = {1: "Initial", 2: "Basic", 3: "Intermediate", 4: "Advanced", 5: "Optimizing"}


def _timestamp() -> str:
    return datetime.now().strftime("%Y-%m-%d_%H%M%S")


def _score_avg(items: list) -> Optional[float]:
    scored = [i for i in items if i.get("score") is not None and not i.get("na", False)]
    if not scored:
        return None
    return sum(i["score"] for i in scored) / len(scored)


def _get_maturity_band(score: float) -> dict:
    for band in MATURITY_BANDS:
        if band["min"] <= score < band["max"]:
            return band
    if score >= 5.0:
        return MATURITY_BANDS[-1]
    return MATURITY_BANDS[0]


def _weighted_composite(d: dict) -> Optional[float]:
    weights = d.get("scoring_config", {}).get("cc_weights", {})
    total_weight = 0.0
    weighted_sum = 0.0
    for mc in d.get("macro_capabilities", []):
        for cc in mc.get("critical_capabilities", []):
            items = []
            for ca in cc.get("capability_areas", []):
                items.extend(ca.get("items", []))
            score = _score_avg(items)
            weight = weights.get(cc["id"], cc.get("weight", 0.028))
            if score is not None:
                weighted_sum += score * weight
                total_weight += weight
    if total_weight == 0:
        return None
    return weighted_sum / total_weight


class ExportEngine:
    def __init__(self, base_dir: str, resource_dir: str | None = None):
        self.base_dir = Path(base_dir)
        self.resource_dir = Path(resource_dir) if resource_dir else self.base_dir
        self.exports_dir = self.base_dir / "exports"
        self.templates_dir = self.resource_dir / "templates"

    def _ensure_exports_dir(self):
        self.exports_dir.mkdir(exist_ok=True)

    def _data_dict(self, data: AssessmentData) -> dict:
        return json.loads(data.model_dump_json())

    def generate_radar_chart_png(self, data: AssessmentData) -> str:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

        d = self._data_dict(data)
        labels = []
        scores = []
        targets = []
        target_scores = d.get("target_scores", {})
        for mc in d.get("macro_capabilities", []):
            for cc in mc.get("critical_capabilities", []):
                labels.append(cc["name"])
                items = []
                for ca in cc.get("capability_areas", []):
                    items.extend(ca.get("items", []))
                s = _score_avg(items)
                scores.append(s if s is not None else 0)
                targets.append(target_scores.get(cc["id"], 3.0))

        n = len(labels)
        if n == 0:
            fig, ax = plt.subplots(figsize=(8, 8))
            ax.text(0.5, 0.5, "No data", ha="center", va="center")
            self._ensure_exports_dir()
            path = str(self.exports_dir / "radar_chart.png")
            fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
            plt.close(fig)
            return path

        angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
        scores_plot = scores + [scores[0]]
        targets_plot = targets + [targets[0]]
        angles_plot = angles + [angles[0]]

        fig, ax = plt.subplots(figsize=(10, 10), subplot_kw=dict(polar=True))
        ax.fill(angles_plot, scores_plot, alpha=0.2, color="#1BA1E2")
        ax.plot(angles_plot, scores_plot, color="#1BA1E2", linewidth=2, label="Current")
        ax.plot(angles_plot, targets_plot, color="#8A8A8E", linewidth=1.5, linestyle="--", label="Target")
        ax.set_xticks(angles)
        ax.set_xticklabels(labels, size=6)
        ax.set_ylim(0, 5)
        ax.set_yticks([1, 2, 3, 4, 5])
        ax.set_yticklabels(["1", "2", "3", "4", "5"], size=8)
        ax.set_title("IT-CMF Maturity Profile", size=14, pad=20)
        ax.legend(loc="upper right", bbox_to_anchor=(1.3, 1.1), fontsize=9)

        self._ensure_exports_dir()
        path = str(self.exports_dir / "radar_chart.png")
        fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return path

    def export_findings(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"ITCMF_Assessment_Findings_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        from docx import Document
        doc = Document()
        doc.add_heading("IT-CMF Assessment Findings", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Industry: {info.get('industry', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")
        doc.add_paragraph(f"Assessor: {info.get('assessor', '')}")

        composite = _weighted_composite(d)
        if composite:
            band = _get_maturity_band(composite)
            doc.add_heading("Overall Maturity", level=1)
            doc.add_paragraph(f"Composite Score: {composite:.2f} — {band['label']}")

        for mc in d.get("macro_capabilities", []):
            doc.add_heading(mc["name"], level=1)
            for cc in mc.get("critical_capabilities", []):
                items = []
                for ca in cc.get("capability_areas", []):
                    items.extend(ca.get("items", []))
                score = _score_avg(items)
                doc.add_heading(cc["name"], level=2)
                doc.add_paragraph(f"Score: {score:.2f} — {_get_maturity_band(score)['label']}" if score else "Score: Not yet scored")

                for ca in cc.get("capability_areas", []):
                    ca_score = _score_avg(ca.get("items", []))
                    doc.add_heading(ca["name"], level=3)
                    doc.add_paragraph(f"Average Score: {ca_score:.2f}" if ca_score else "Not scored")

                    for item in ca.get("items", []):
                        score_val = item.get("score")
                        na = item.get("na", False)
                        if na:
                            text = f"[N/A] {item['text']}"
                        elif score_val:
                            text = f"[{score_val} - {SCORE_LABELS.get(score_val, '')}] {item['text']}"
                        else:
                            text = f"[--] {item['text']}"
                        if item.get("notes"):
                            text += f"\n  Notes: {item['notes']}"
                        doc.add_paragraph(text, style="List Bullet")

        doc.save(str(output_path))
        return filename

    def export_executive_summary(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"ITCMF_Executive_Summary_{_timestamp()}.docx"
        output_path = self.exports_dir / filename
        chart_path = self.generate_radar_chart_png(data)

        from docx import Document
        from docx.shared import Inches
        doc = Document()
        doc.add_heading("Executive Summary — IT-CMF Assessment", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Industry: {info.get('industry', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        composite = _weighted_composite(d)
        if composite:
            band = _get_maturity_band(composite)
            doc.add_heading("Overall Maturity", level=1)
            doc.add_paragraph(f"Score: {composite:.2f} — {band['label']}")

        doc.add_heading("Maturity Profile", level=1)
        doc.add_picture(chart_path, width=Inches(5))

        doc.add_heading("Critical Capability Scores", level=1)
        for mc in d.get("macro_capabilities", []):
            doc.add_heading(mc["name"], level=2)
            for cc in mc.get("critical_capabilities", []):
                items = []
                for ca in cc.get("capability_areas", []):
                    items.extend(ca.get("items", []))
                score = _score_avg(items)
                doc.add_paragraph(
                    f"{cc['name']}: {score:.2f} — {_get_maturity_band(score)['label']}" if score else f"{cc['name']}: N/A"
                )

        # Top gaps
        target_scores = d.get("target_scores", {})
        gaps = []
        for mc in d.get("macro_capabilities", []):
            for cc in mc.get("critical_capabilities", []):
                items = []
                for ca in cc.get("capability_areas", []):
                    items.extend(ca.get("items", []))
                current = _score_avg(items)
                target = target_scores.get(cc["id"], 3.0)
                if current is not None:
                    gap = target - current
                    if gap > 0:
                        gaps.append({"name": cc["name"], "current": current, "target": target, "gap": gap})
        gaps.sort(key=lambda g: g["gap"], reverse=True)
        if gaps:
            doc.add_heading("Top Priority Gaps", level=1)
            for g in gaps[:5]:
                doc.add_paragraph(f"{g['name']}: Current {g['current']:.2f} → Target {g['target']:.1f} (Gap: {g['gap']:.2f})")

        doc.save(str(output_path))
        return filename

    def export_gap_analysis(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"ITCMF_Gap_Analysis_{_timestamp()}.docx"
        output_path = self.exports_dir / filename
        target_scores = d.get("target_scores", {})

        from docx import Document
        doc = Document()
        doc.add_heading("Gap Analysis & Roadmap", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        doc.add_heading("Gap Matrix", level=1)
        table = doc.add_table(rows=1, cols=5)
        table.style = "Table Grid"
        for i, h in enumerate(["Critical Capability", "Current", "Target", "Gap", "Severity"]):
            table.rows[0].cells[i].text = h

        for mc in d.get("macro_capabilities", []):
            for cc in mc.get("critical_capabilities", []):
                items = []
                for ca in cc.get("capability_areas", []):
                    items.extend(ca.get("items", []))
                current = _score_avg(items)
                target = target_scores.get(cc["id"], 3.0)
                gap = (target - current) if current else None
                severity = "High" if gap and gap > 1.5 else "Medium" if gap and gap > 0.5 else "Low"
                row = table.add_row().cells
                row[0].text = cc["name"]
                row[1].text = f"{current:.2f}" if current else "N/A"
                row[2].text = f"{target:.1f}"
                row[3].text = f"{gap:.2f}" if gap is not None else "N/A"
                row[4].text = severity

        doc.add_heading("Remediation Roadmap", level=1)
        doc.add_heading("30-Day Quick Wins", level=2)
        doc.add_paragraph("Focus on capabilities with scores at Initial (1) for immediate improvement opportunities.")
        doc.add_heading("60-Day Improvements", level=2)
        doc.add_paragraph("Establish formal processes for capabilities scoring Basic (2).")
        doc.add_heading("90-Day Milestones", level=2)
        doc.add_paragraph("Target Intermediate maturity (3.0+) for critical capabilities.")
        doc.add_heading("6-12 Month Goals", level=2)
        doc.add_paragraph("Achieve Advanced or Optimizing maturity across all capabilities with continuous improvement cycles.")

        doc.save(str(output_path))
        return filename

    def export_workbook(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"ITCMF_Scored_Workbook_{_timestamp()}.xlsx"
        output_path = self.exports_dir / filename

        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill

        wb = Workbook()
        ws = wb.active
        ws.title = "Dashboard"
        ws["A1"] = "IT-CMF Assessment Dashboard"
        ws["A1"].font = Font(bold=True, size=14)
        ws["A3"] = "Client:"
        ws["B3"] = d.get("client_info", {}).get("name", "")
        ws["A4"] = "Industry:"
        ws["B4"] = d.get("client_info", {}).get("industry", "")
        ws["A5"] = "Assessment Date:"
        ws["B5"] = d.get("client_info", {}).get("assessment_date", "")
        ws["A6"] = "Assessor:"
        ws["B6"] = d.get("client_info", {}).get("assessor", "")

        composite = _weighted_composite(d)
        ws["A8"] = "Overall Maturity Score:"
        ws["B8"] = round(composite, 2) if composite else "N/A"
        if composite:
            ws["C8"] = _get_maturity_band(composite)["label"]

        row = 10
        header_fill = PatternFill(start_color="1BA1E2", end_color="1BA1E2", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF")
        for col, h in enumerate(["Macro-Capability", "Critical Capability", "Weight", "Score", "Band"], 1):
            ws.cell(row=row, column=col, value=h).font = header_font
            ws.cell(row=row, column=col).fill = header_fill

        for mc in d.get("macro_capabilities", []):
            for cc in mc.get("critical_capabilities", []):
                row += 1
                items = []
                for ca in cc.get("capability_areas", []):
                    items.extend(ca.get("items", []))
                score = _score_avg(items)
                ws.cell(row=row, column=1, value=mc["name"])
                ws.cell(row=row, column=2, value=cc["name"])
                ws.cell(row=row, column=3, value=f"{cc.get('weight', 0) * 100:.1f}%")
                ws.cell(row=row, column=4, value=round(score, 2) if score else "N/A")
                if score:
                    ws.cell(row=row, column=5, value=_get_maturity_band(score)["label"])

        # Per-CC sheets
        for mc in d.get("macro_capabilities", []):
            for cc in mc.get("critical_capabilities", []):
                sheet_name = cc["name"][:31]
                ws = wb.create_sheet(title=sheet_name)
                ws["A1"] = cc["name"]
                ws["A1"].font = Font(bold=True, size=12)
                ws["A2"] = f"Macro-Capability: {mc['name']}"
                ws["A2"].font = Font(italic=True, color="666666")

                row = 4
                headers = ["Capability Area", "Item #", "Assessment Item", "Score", "Confidence", "Evidence", "Notes"]
                for col, h in enumerate(headers, 1):
                    cell = ws.cell(row=row, column=col, value=h)
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill(start_color="E2E8F0", end_color="E2E8F0", fill_type="solid")

                for ca in cc.get("capability_areas", []):
                    for item in ca.get("items", []):
                        row += 1
                        ws.cell(row=row, column=1, value=ca["name"])
                        ws.cell(row=row, column=2, value=item["id"])
                        ws.cell(row=row, column=3, value=item["text"])
                        ws.cell(row=row, column=4, value=item.get("score"))
                        ws.cell(row=row, column=5, value=item.get("confidence", ""))
                        refs = item.get("evidence_references", [])
                        evidence_str = "; ".join(
                            f"{r.get('document', '')} §{r.get('section', '')}"
                            for r in refs if r.get("document")
                        ) if refs else ""
                        ws.cell(row=row, column=6, value=evidence_str)
                        ws.cell(row=row, column=7, value=item.get("notes", ""))

                for col in ws.columns:
                    max_len = 0
                    col_letter = col[0].column_letter
                    for cell in col:
                        if cell.value:
                            max_len = max(max_len, len(str(cell.value)))
                    ws.column_dimensions[col_letter].width = min(max_len + 2, 60)

        wb.save(str(output_path))
        return filename

    def export_outbrief(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"ITCMF_OutBrief_{_timestamp()}.pptx"
        output_path = self.exports_dir / filename
        chart_path = self.generate_radar_chart_png(data)

        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "IT-CMF Assessment Out-Brief"
        subtitle = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
        if subtitle:
            subtitle.text = d.get("client_info", {}).get("name", "")

        # Overview slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Assessment Overview"
        composite = _weighted_composite(d)
        body = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
        if body:
            tf = body.text_frame
            tf.text = f"Client: {d.get('client_info', {}).get('name', '')}"
            tf.add_paragraph().text = f"Industry: {d.get('client_info', {}).get('industry', '')}"
            tf.add_paragraph().text = f"Date: {d.get('client_info', {}).get('assessment_date', '')}"
            tf.add_paragraph().text = f"Assessor: {d.get('client_info', {}).get('assessor', '')}"
            if composite:
                band = _get_maturity_band(composite)
                tf.add_paragraph().text = f"Overall Score: {composite:.2f} — {band['label']}"

        # Radar chart slide
        blank_idx = min(5, len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[blank_idx])
        slide.shapes.add_picture(chart_path, Inches(3), Inches(0.5), Inches(7), Inches(6.5))

        # Per-macro-capability slides
        for mc in d.get("macro_capabilities", []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = mc["name"]
            body = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
            if body:
                tf = body.text_frame
                tf.text = ""
                for cc in mc.get("critical_capabilities", []):
                    items = []
                    for ca in cc.get("capability_areas", []):
                        items.extend(ca.get("items", []))
                    score = _score_avg(items)
                    p = tf.add_paragraph()
                    p.text = f"{cc['name']}: {score:.2f} — {_get_maturity_band(score)['label']}" if score else f"{cc['name']}: Not scored"

        prs.save(str(output_path))
        return filename

    def export_heatmap(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"ITCMF_Maturity_Heatmap_{_timestamp()}.xlsx"
        output_path = self.exports_dir / filename

        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill

        wb = Workbook()
        ws = wb.active
        ws.title = "Heatmap"

        ws["A1"] = "IT-CMF Maturity Heatmap"
        ws["A1"].font = Font(bold=True, size=14)

        score_colors = {
            1: "EF4444", 2: "F97316", 3: "EAB308", 4: "84CC16", 5: "22C55E",
        }

        row = 3
        ws.cell(row=row, column=1, value="Critical Capability").font = Font(bold=True)

        max_cas = 0
        for mc in d.get("macro_capabilities", []):
            for cc in mc.get("critical_capabilities", []):
                max_cas = max(max_cas, len(cc.get("capability_areas", [])))

        for i in range(max_cas):
            ws.cell(row=row, column=i + 2, value=f"CA {i + 1}").font = Font(bold=True)
        ws.cell(row=row, column=max_cas + 2, value="CC Avg").font = Font(bold=True)

        for mc in d.get("macro_capabilities", []):
            for cc in mc.get("critical_capabilities", []):
                row += 1
                ws.cell(row=row, column=1, value=cc["name"])
                all_items = []
                for ca_idx, ca in enumerate(cc.get("capability_areas", [])):
                    all_items.extend(ca.get("items", []))
                    ca_score = _score_avg(ca.get("items", []))
                    cell = ws.cell(row=row, column=ca_idx + 2)
                    if ca_score is not None:
                        cell.value = round(ca_score, 2)
                        rounded = min(5, max(1, round(ca_score)))
                        color = score_colors.get(rounded, "FFFFFF")
                        cell.fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
                        if rounded <= 2:
                            cell.font = Font(color="FFFFFF", bold=True)
                        else:
                            cell.font = Font(bold=True)
                    else:
                        cell.value = "--"
                cc_avg = _score_avg(all_items)
                avg_cell = ws.cell(row=row, column=max_cas + 2)
                if cc_avg is not None:
                    avg_cell.value = round(cc_avg, 2)
                    avg_cell.font = Font(bold=True)

        row += 2
        ws.cell(row=row, column=1, value="Legend:").font = Font(bold=True)
        for score, label in SCORE_LABELS.items():
            row += 1
            cell = ws.cell(row=row, column=1, value=f"{score} — {label}")
            color = score_colors[score]
            cell.fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
            if score <= 2:
                cell.font = Font(color="FFFFFF")

        ws.column_dimensions["A"].width = 35
        for i in range(max_cas + 2):
            ws.column_dimensions[chr(66 + i)].width = 12

        wb.save(str(output_path))
        return filename

    def export_quick_wins(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"ITCMF_Quick_Wins_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        from docx import Document
        doc = Document()
        doc.add_heading("IT-CMF Quick Wins Report", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        doc.add_heading("Methodology", level=1)
        doc.add_paragraph(
            "Quick wins are identified as assessment items scored 1-2 (Initial or Basic) "
            "that have the highest potential for improvement. Items are prioritized by CC weight "
            "multiplied by the score gap to target."
        )

        target_scores = d.get("target_scores", {})
        weights = d.get("scoring_config", {}).get("cc_weights", {})
        candidates = []
        for mc in d.get("macro_capabilities", []):
            for cc in mc.get("critical_capabilities", []):
                weight = weights.get(cc["id"], 0.028)
                target = target_scores.get(cc["id"], 3.0)
                for ca in cc.get("capability_areas", []):
                    for item in ca.get("items", []):
                        score = item.get("score")
                        if score is not None and score <= 2 and not item.get("na", False):
                            gap = target - score
                            priority = weight * gap
                            candidates.append({
                                "cc": cc["name"],
                                "ca": ca["name"],
                                "item_id": item["id"],
                                "text": item["text"],
                                "score": score,
                                "target": target,
                                "gap": gap,
                                "priority": priority,
                                "notes": item.get("notes", ""),
                            })

        candidates.sort(key=lambda c: c["priority"], reverse=True)

        if not candidates:
            doc.add_heading("Results", level=1)
            doc.add_paragraph("No quick win candidates found. All items are scored 3 or above, or are unscored.")
        else:
            from collections import defaultdict
            by_cc: dict[str, list] = defaultdict(list)
            for c in candidates[:30]:
                by_cc[c["cc"]].append(c)

            for cc_name, items in by_cc.items():
                doc.add_heading(cc_name, level=1)
                table = doc.add_table(rows=1, cols=4)
                table.style = "Table Grid"
                for i, h in enumerate(["Item", "Current", "Target", "Gap"]):
                    table.rows[0].cells[i].text = h
                for item in items:
                    row = table.add_row().cells
                    row[0].text = f"[{item['item_id']}] {item['text'][:80]}"
                    row[1].text = f"{item['score']} ({SCORE_LABELS.get(item['score'], '')})"
                    row[2].text = f"{item['target']:.1f}"
                    row[3].text = f"{item['gap']:.1f}"

        doc.save(str(output_path))
        return filename

    def export_cc_roadmap(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"ITCMF_CC_Roadmap_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        from docx import Document
        doc = Document()
        doc.add_heading("CC Improvement Roadmap", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        target_scores = d.get("target_scores", {})

        for mc in d.get("macro_capabilities", []):
            doc.add_heading(mc["name"], level=1)
            for cc in mc.get("critical_capabilities", []):
                items = []
                for ca in cc.get("capability_areas", []):
                    items.extend(ca.get("items", []))
                score = _score_avg(items)
                target = target_scores.get(cc["id"], 3.0)
                gap = (target - score) if score else None

                doc.add_heading(cc["name"], level=2)
                if score:
                    doc.add_paragraph(f"Current: {score:.2f} ({_get_maturity_band(score)['label']}) → Target: {target:.1f}")
                    if gap and gap > 0:
                        doc.add_paragraph(f"Gap: {gap:.2f} — Priority: {'High' if gap > 1.5 else 'Medium' if gap > 0.5 else 'Low'}")
                else:
                    doc.add_paragraph("Not yet scored")

                # List low-scoring areas
                for ca in cc.get("capability_areas", []):
                    ca_score = _score_avg(ca.get("items", []))
                    if ca_score is not None and ca_score < target:
                        doc.add_paragraph(
                            f"  {ca['name']}: {ca_score:.2f} — needs improvement",
                            style="List Bullet"
                        )

        doc.save(str(output_path))
        return filename

    def export_all(self, data: AssessmentData) -> list[str]:
        return [
            self.export_findings(data),
            self.export_executive_summary(data),
            self.export_gap_analysis(data),
            self.export_workbook(data),
            self.export_outbrief(data),
            self.export_heatmap(data),
            self.export_quick_wins(data),
            self.export_cc_roadmap(data),
        ]
